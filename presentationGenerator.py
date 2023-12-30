import requests
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR

def download_image(url, filename):
    response = requests.get(url)
    with open(filename, "wb") as f:
        f.write(response.content)

def generate_slideshow(excel_file):
    # Load the Excel workbook
    workbook = load_workbook(excel_file)
    sheet = workbook.active

    # Create a PowerPoint presentation
    presentation = Presentation()

    # Set slide dimensions
    slide_width = Inches(13.33)
    slide_height = Inches(7.5)
    presentation.slide_width = slide_width
    presentation.slide_height = slide_height

    # Set font properties
    font_name = "TW Cen MT"
    font_size = Pt(18)
    alt_font_size = Pt(20)

    # Create a blue gradient background
    slide_master = presentation.slide_master
    slide_master.background.fill.gradient()
    slide_master.background.fill.gradient_angle = 90
    slide_master.background.fill.gradient_stops[0].position = 0
    slide_master.background.fill.gradient_stops[0].color.rgb = RGBColor(0x93, 0xC4, 0xEA)
    slide_master.background.fill.gradient_stops[1].position = 1
    slide_master.background.fill.gradient_stops[1].color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Iterate through each row in the sheet
    for row in sheet.iter_rows(min_row=2, values_only=True):
        # Create a new slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        # Set slide dimensions
        slide.slide_width = slide_width
        slide.slide_height = slide_height

        # Create a table on the slide
        table_width = Inches(4.17)
        table_height = slide_height
        left = 0  # Align the left border of the table with the left border of the slide
        top = 0  # Align the top border of the table with the top border of the slide

        # Determine the number of columns and rows
        column_order = [0, 8, 4, 5, 2, 3, 7, 9]  # A, I, E, F, C, D, H, J
        num_columns = 2
        cell_texts = ["Name", "Height (Feet)", "Height (Meters)", "Floors", "Status", "Completion", "Function", "Design"]
        num_rows = 8

        table = slide.shapes.add_table(
            num_rows, num_columns, left, top, table_width, table_height
        ).table
        # Set table style to medium style 3 accent 1
        table.style = "Table Grid"

        # Calculate the row height based on the table height
        row_height = table_height / num_rows

        # Set column widths
        column_widths = [Inches(1.45), Inches(2.73)]
        for idx, column_width in enumerate(column_widths):
            table.columns[idx].width = column_width

        # Manually fill the cells in column 0
        for idx, cell_text in enumerate(cell_texts):
            cell = table.cell(idx, 0)
            cell.text = cell_text
            # Center align the text
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # Set font properties
            cell.text_frame.paragraphs[0].font.name = font_name
            cell.text_frame.paragraphs[0].font.size = font_size

        # Manually set text color and bold for cells in column 0
        for row_idx in range(1, num_rows):
            cell = table.cell(row_idx, 0)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White
            cell.text_frame.paragraphs[0].font.bold = True

        # Manually set shading for cells in column 0
        for row_idx in range(num_rows):
            cell = table.cell(row_idx, 0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)  # Blue, Accent 1, Darker 25%
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


        # Populate the table with data from the current row
        for idx, col_idx in enumerate(column_order):
            cell = table.cell(idx, 1)
            cell_value = str(row[col_idx])
            if idx == 1:  # Check if it's the second column, second cell
                cell_value += " FT"
            if idx == 2:  # Check if it's the second column, second cell
                cell_value += " M"
            cell.text = cell_value
            # Center align the text
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Set font properties
            cell.text_frame.paragraphs[0].font.name = font_name
            cell.text_frame.paragraphs[0].font.size = alt_font_size

            # Manually set shading for cell (0, 1)
            cell = table.cell(0, 1)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)  # Blue, Accent 1, Darker 25%
            # Center align the text
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # Set font properties
            cell.text_frame.paragraphs[0].font.name = font_name
            cell.text_frame.paragraphs[0].font.size = font_size
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            #Setting bottom right cell to smaller text
            cell = table.cell(7,1)
            cell.text_frame.paragraphs[0].font.size = font_size

        # Check if image URL is "No image found"
        image_url = row[-1]
        if image_url.lower() != "no image found":
            # Download the image and save it locally
            image_filename = "image.jpg"
            download_image(image_url, image_filename)

            # Add the downloaded image to the right of the table
            image = slide.shapes.add_picture(
                image_filename, left + table_width, top, height=slide_height
            )

            # Scale the image width proportionally
            available_width = slide_width - table_width
            aspect_ratio = image.width / image.height
            image.width = int(slide_height * aspect_ratio)
            image.height = slide_height
            # Shift the image to the right
            image.left += Inches(1.7)

    # Save the PowerPoint presentation
    output_file = "slideshow.pptx"
    presentation.save(output_file)
    print("Slideshow generated successfully. Saved as", output_file)

# Provide the path to your Excel file
excel_file = "combined_data_second.xlsx"

# Generate the slideshow
generate_slideshow(excel_file)